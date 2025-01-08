import os, sys
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

#pdf_file = 'in.pdf'
# Prep presentation

def pdf_to_power(pdf_file):
	prs = Presentation()
	blank_slide_layout = prs.slide_layouts[6]
	base_name = pdf_file.split(".pdf")[0]

	# Convert PDF to list of images
	slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2,poppler_path=r'.\poppler\bin')

	# Loop over slides
	for i, slideimg in enumerate(slideimgs):
		imagefile = BytesIO()
		slideimg.save(imagefile, format='tiff')
		imagedata = imagefile.getvalue()
		imagefile.seek(0)
		width, height = slideimg.size
		prs.slide_height = height * 9525
		prs.slide_width = width * 9525
		slide = prs.slides.add_slide(blank_slide_layout)
		pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

	# Save Powerpoint
	prs.save(base_name + '.pptx')
